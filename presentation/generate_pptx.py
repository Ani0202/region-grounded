"""
Generate GRAFT presentation (v2) as a PowerPoint file.
Usage: python generate_pptx.py
Output: graft_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ────────────────────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
BASE_DIR   = os.path.dirname(SCRIPT_DIR)

IMG_VIZ_01   = os.path.join(BASE_DIR, "report", "pg_viz_01_1836335410.png")
IMG_VIZ_04   = os.path.join(BASE_DIR, "report", "pg_viz_04_1507563902.png")
IMG_FLICKR   = os.path.join(BASE_DIR, "notebooks", "outputs", "flickr30k_entities_grid.png")
IMG_FLORENCE = os.path.join(BASE_DIR, "notebooks", "florence2_viz.png")

# ── Colour palette ────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD    = RGBColor(0x1E, 0x29, 0x3B)
BG_DARK2   = RGBColor(0x17, 0x22, 0x34)
ACCENT     = RGBColor(0x38, 0xBD, 0xF8)
ACCENT2    = RGBColor(0x81, 0x8C, 0xF8)
GREEN      = RGBColor(0x4A, 0xDE, 0x80)
RED        = RGBColor(0xF8, 0x71, 0x71)
AMBER      = RGBColor(0xFB, 0xBF, 0x24)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
MED_GRAY   = RGBColor(0x64, 0x74, 0x8B)
BORDER     = RGBColor(0x33, 0x41, 0x55)
HDR_BLUE   = RGBColor(0x1E, 0x40, 0x6E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, border_color=None, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_text(slide, left, top, width, height, text, size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def add_bullets(slide, left, top, width, height, items, size=16,
                color=WHITE, dot=ACCENT, spacing=Pt(5)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        rb = p.add_run()
        rb.text = "▸ "
        rb.font.size = Pt(size)
        rb.font.color.rgb = dot
        rb.font.name = "Calibri"
        if isinstance(item, tuple):
            r1 = p.add_run()
            r1.text = item[0]
            r1.font.size = Pt(size)
            r1.font.color.rgb = color
            r1.font.bold = True
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = item[1]
            r2.font.size = Pt(size)
            r2.font.color.rgb = LIGHT_GRAY
            r2.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def add_line(slide, left, top, width, color=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    ts = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    tbl = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(13); p.font.bold = True
            p.font.color.rgb = WHITE; p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        c.fill.solid(); c.fill.fore_color.rgb = HDR_BLUE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            if isinstance(val, tuple):
                c.text = val[0]; txt_col = val[1]
            else:
                c.text = str(val); txt_col = WHITE
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(12); p.font.color.rgb = txt_col
                p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER
            c.fill.solid()
            c.fill.fore_color.rgb = BG_CARD if i % 2 == 0 else BG_DARK2
    return ts


def badge(slide, text, color=ACCENT, w=Inches(2.8)):
    add_rect(slide, Inches(0.8), Inches(0.4), w, Inches(0.45), color, corner_radius=0.15)
    add_text(slide, Inches(0.8), Inches(0.4), w, Inches(0.45),
             text, size=16, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)


def slide_title(slide, text, size=32, accent_w=None):
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.7),
             text, size=size, color=ACCENT, bold=True)
    w = accent_w or Inches(min(len(text) * 0.19 + 0.4, 11.5))
    add_line(slide, Inches(0.8), Inches(1.1), w)


def bold_normal_row(slide, left, top, width, height, bold_text, normal_text, size=17):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = bold_text
    r1.font.size = Pt(size); r1.font.color.rgb = WHITE
    r1.font.bold = True; r1.font.name = "Calibri"
    r2 = p.add_run(); r2.text = normal_text
    r2.font.size = Pt(size); r2.font.color.rgb = LIGHT_GRAY
    r2.font.name = "Calibri"


# ═════════════════════════════════════════════════════════════════════════════
# BUILD
# ═════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BL = prs.slide_layouts[6]   # blank


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1  Title
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
add_text(sl, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "GRAFT", size=54, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
         "Grounding Region Annotations into Fine-Tuned Transformers",
         size=26, color=WHITE, align=PP_ALIGN.CENTER)
add_line(sl, Inches(4.5), Inches(3.55), Inches(4.3))
add_text(sl, Inches(1), Inches(3.9), Inches(11), Inches(0.5),
         "TTIC 31280 — Models for Representing Images & Videos",
         size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(4.65), Inches(11), Inches(0.5),
         "Siddharth Raj  ·  Aniket Agrawal",
         size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(5.25), Inches(11), Inches(0.5),
         "University of Chicago  |  Spring 2026",
         size=16, color=MED_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2  Outline
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Outline", size=36, accent_w=Inches(2.5))

sections = [
    ("1", "Motivation & Setup",    "~5 min", "The problem, architecture, data, evaluation"),
    ("2", "Main Findings",         "~5 min", "Representation mismatches, variance, ablations"),
    ("3", "Results & Surprises",   "~4 min", "Qualitative results, staircase training, auto-labels"),
    ("4", "Future Work",           "~1 min", "Next steps and open questions"),
]
for i, (num, title, time, desc) in enumerate(sections):
    y = Inches(1.55 + i * 1.3)
    add_rect(sl, Inches(1.2), y, Inches(0.6), Inches(0.6), ACCENT, corner_radius=0.5)
    add_text(sl, Inches(1.2), y + Pt(3), Inches(0.6), Inches(0.6),
             num, size=22, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, Inches(2.2), y, Inches(7), Inches(0.38),
             title, size=22, color=WHITE, bold=True)
    add_text(sl, Inches(2.2), y + Inches(0.38), Inches(7), Inches(0.35),
             desc, size=14, color=LIGHT_GRAY)
    add_rect(sl, Inches(10.0), y + Inches(0.05), Inches(1.3), Inches(0.42),
             BG_CARD, border_color=BORDER, corner_radius=0.15)
    add_text(sl, Inches(10.0), y + Inches(0.05), Inches(1.3), Inches(0.42),
             time, size=14, color=ACCENT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3  The Problem  (with Flickr30k entities grid)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "The Problem: Spatial Blindness in VLMs", size=30, accent_w=Inches(6))

add_bullets(sl, Inches(0.8), Inches(1.35), Inches(5.8), Inches(2.3), [
    ("CLIP / SigLIP ", "train with global image–text contrastive objectives"),
    ("Patch features are spatially coherent, ", "but phrase-to-patch alignment is weak"),
    ("This limits: ", "referring expression comprehension, dense captioning, spatial QA"),
], size=16)

add_rect(sl, Inches(0.8), Inches(3.8), Inches(5.8), Inches(2.8),
         BG_CARD, border_color=ACCENT2, corner_radius=0.03)
add_text(sl, Inches(1.1), Inches(4.0), Inches(5.3), Inches(0.4),
         "Why SigLIP?", size=20, color=ACCENT2, bold=True)
add_bullets(sl, Inches(1.1), Inches(4.5), Inches(5.3), Inches(1.9), [
    ("Sigmoid loss — ", "no batch denominator → stronger image-text matching"),
    ("Less spatial bias than CLIP ", "— a harder, more honest test case"),
    ("Success here ", "implies generalisation to other contrastive VLMs"),
], size=14, dot=ACCENT2)

# Flickr30k entities grid  (2349×2368, ~square)  width=5.6" → height≈5.64"
sl.shapes.add_picture(IMG_FLICKR, Inches(7.1), Inches(1.1), width=Inches(5.6))
add_text(sl, Inches(7.1), Inches(6.8), Inches(5.6), Inches(0.35),
         "Flickr30k Entities — phrase → bounding box annotations",
         size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4  GRAFT Architecture
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "GRAFT: What We Built", accent_w=Inches(4.2))

add_rect(sl, Inches(1.8), Inches(1.35), Inches(9.7), Inches(0.72),
         BG_CARD, border_color=ACCENT, corner_radius=0.03)
add_text(sl, Inches(1.8), Inches(1.4), Inches(9.7), Inches(0.72),
         "Freeze SigLIP  +  LoRA adapters  +  Region loss   →   ℒ = ℒ_global + 0.5 · ℒ_region",
         size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(0.8), Inches(2.35), Inches(5.6), Inches(4.35),
         BG_CARD, border_color=BORDER, corner_radius=0.03)
add_text(sl, Inches(1.1), Inches(2.55), Inches(5.1), Inches(0.4),
         "Architecture", size=22, color=GREEN, bold=True)
add_bullets(sl, Inches(1.1), Inches(3.05), Inches(5.1), Inches(3.5), [
    ("LoRA on q-proj, v-proj ", "(both encoders)"),
    ("Rank 4, α=16 ", "→ 295K trainable params (0.14% of backbone)"),
    ("MaskCLIP bypass ", "installed before training — same path as eval"),
    ("Best-PG snapshot: ", "checkpoint LoRA weights every 200 steps (~1 MB)"),
], size=15, dot=GREEN)

add_rect(sl, Inches(6.9), Inches(2.35), Inches(5.7), Inches(4.35),
         BG_CARD, border_color=BORDER, corner_radius=0.03)
add_text(sl, Inches(7.2), Inches(2.55), Inches(5.2), Inches(0.4),
         "Training Losses", size=22, color=AMBER, bold=True)
add_bullets(sl, Inches(7.2), Inches(3.05), Inches(5.2), Inches(3.5), [
    ("ℒ_global: ", "sigmoid-SigLIP on MHAP pooler ↔ EOS caption embedding"),
    ("ℒ_region: ", "FILIP-style EOS phrase ↔ patch max-pool over in-bbox patches"),
    ("λ = 0.5: ", "region weight — higher destabilises training"),
    ("Optimiser: ", "AdamW, lr=2e-4, cosine decay, bf16 autocast"),
], size=15, dot=AMBER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5  Data & Evaluation
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Data & Evaluation", accent_w=Inches(3.8))

add_rect(sl, Inches(0.8), Inches(1.5), Inches(5.6), Inches(3.1),
         BG_CARD, border_color=BORDER, corner_radius=0.03)
add_text(sl, Inches(1.1), Inches(1.7), Inches(5.1), Inches(0.4),
         "Training Data", size=20, color=GREEN, bold=True)
add_bullets(sl, Inches(1.1), Inches(2.2), Inches(5.1), Inches(2.2), [
    ("Flickr30k-Entities ", "(29K train images)"),
    ("Human-annotated phrase → bbox pairs ", "(~5.7 phrases/image)"),
    ("Also tested: ", "Florence-2 auto-annotations (two modes)"),
    ("15 experiments total ", "iterating to the locked recipe"),
], size=15, dot=GREEN)

add_rect(sl, Inches(6.9), Inches(1.5), Inches(5.7), Inches(3.1),
         BG_CARD, border_color=BORDER, corner_radius=0.03)
add_text(sl, Inches(7.2), Inches(1.7), Inches(5.2), Inches(0.4),
         "Evaluation Metrics", size=20, color=ACCENT2, bold=True)
add_bullets(sl, Inches(7.2), Inches(2.2), Inches(5.2), Inches(2.2), [
    ("Pointing Game (PG): ", "argmax patch centre inside GT bbox?"),
    ("Recall@1: ", "tight box over top-k patches overlaps GT at IoU≥0.5"),
    ("Val: ", "200 images, 2,124 phrase-bbox pairs (fixed)"),
    ("Test: ", "2,900 held-out images, 28,467 pairs (SE≈0.24 pp)"),
], size=15, dot=ACCENT2)

add_rect(sl, Inches(0.8), Inches(4.95), Inches(11.7), Inches(1.45),
         BG_CARD, border_color=AMBER, corner_radius=0.03)
add_text(sl, Inches(1.1), Inches(5.1), Inches(4.5), Inches(0.4),
         "Frozen Baseline (B0) — no fine-tuning", size=18, color=AMBER, bold=True)
add_text(sl, Inches(1.1), Inches(5.6), Inches(11.0), Inches(0.6),
         "PG = 14.74%     |     R@1 top-25 = 7.91%     |     R@1 halfmax = 7.58%",
         size=20, color=WHITE, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6  Experiment milestones
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "The Journey: 15 Experiments", accent_w=Inches(5.5))

# 5 milestone columns
milestones = [
    ("B0\nBaseline",          "14.74%",       "",                   LIGHT_GRAY, LIGHT_GRAY),
    ("Exp 03\nMaskCLIP\nBypass ★", "22.08%", "+7.34 pp\n(largest gain)", GREEN,  GREEN),
    ("Exp 06\nSeed Control",  "18.68\n±2.57%","True mean\n+3.94 pp", ACCENT,    ACCENT),
    ("Exp 11\nEOS Fix ★",     "20.76\n±0.41%","6× variance↓\n+2.08 pp mean", ACCENT2, ACCENT2),
    ("Exp 13\n5 Epochs ★",    "24.44%",       "Best result\n+9.70 pp", AMBER,   AMBER),
]
BW = Inches(2.2); GAP = Inches(0.13); Y0 = Inches(1.45); BH = Inches(4.7)
X0 = Inches(0.6)

for i, (label, pg, delta, border, tcol) in enumerate(milestones):
    x = X0 + i * (BW + GAP)
    is_star = "★" in label
    add_rect(sl, x, Y0, BW, BH, BG_CARD, border_color=border, corner_radius=0.03)
    add_text(sl, x + Inches(0.12), Y0 + Inches(0.15), BW - Inches(0.24), Inches(1.05),
             label, size=14, color=tcol, bold=is_star, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.08), Y0 + Inches(1.4), BW - Inches(0.16), Inches(0.85),
             pg, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if delta:
        add_text(sl, x + Inches(0.08), Y0 + Inches(2.45), BW - Inches(0.16), Inches(0.75),
                 delta, size=13, color=tcol, align=PP_ALIGN.CENTER)
    if i < 4:
        ax = x + BW + Inches(0.01)
        add_text(sl, ax, Y0 + Inches(2.1), GAP + Inches(0.1), Inches(0.4),
                 "→", size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)

add_text(sl, Inches(0.6), Inches(6.35), Inches(12.2), Inches(0.45),
         "★ = major fix   |   Ablations between ★s all returned to the same mean (~18–21%)",
         size=13, color=MED_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7  Main Findings — Representation alignment
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
badge(sl, "MAIN FINDINGS")

add_text(sl, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.7),
         "Core Insight: Representation Alignment > Loss Design",
         size=28, color=WHITE, bold=True)
add_line(sl, Inches(0.8), Inches(1.8), Inches(9))
add_text(sl, Inches(1.0), Inches(1.95), Inches(11.5), Inches(0.55),
         "The two biggest gains come from fixing train/eval mismatches — not from changing the loss or the model.",
         size=17, color=LIGHT_GRAY)

# Fix 1
add_rect(sl, Inches(0.8), Inches(2.75), Inches(5.8), Inches(3.85),
         BG_CARD, border_color=GREEN, corner_radius=0.03)
add_text(sl, Inches(1.1), Inches(2.95), Inches(5.3), Inches(0.4),
         "Fix #1 — Image Path: MaskCLIP Bypass", size=18, color=GREEN, bold=True)
add_text(sl, Inches(1.1), Inches(3.42), Inches(5.3), Inches(0.55),
         "Training used full self-attention; evaluation read bypass features → gradients sharpened the wrong representation",
         size=13, color=LIGHT_GRAY)
add_text(sl, Inches(1.1), Inches(4.05), Inches(5.3), Inches(0.45),
         "bypass(h) = W_O · W_V · h", size=17, color=ACCENT, bold=True)
add_bullets(sl, Inches(1.1), Inches(4.6), Inches(5.3), Inches(1.8), [
    ("One-line fix: ", "install bypass before training begins"),
    ("Result: +7.34 pp PG ", "— the largest single gain in the project"),
    ("Shape change: ", "broad weak activation → single sharp spike inside bbox"),
], size=14, dot=GREEN)

# Fix 2
add_rect(sl, Inches(7.1), Inches(2.75), Inches(5.5), Inches(3.85),
         BG_CARD, border_color=ACCENT2, corner_radius=0.03)
add_text(sl, Inches(7.4), Inches(2.95), Inches(5.0), Inches(0.4),
         "Fix #2 — Text Path: EOS Token Alignment", size=18, color=ACCENT2, bold=True)
add_text(sl, Inches(7.4), Inches(3.42), Inches(5.0), Inches(0.55),
         "Region loss used mean(hidden_states); evaluation read only the EOS token → two different vectors",
         size=13, color=LIGHT_GRAY)
add_text(sl, Inches(7.4), Inches(4.05), Inches(5.0), Inches(0.45),
         "e_phrase = last_hidden_state[:, -1, :]", size=17, color=ACCENT2, bold=True)
add_bullets(sl, Inches(7.4), Inches(4.6), Inches(5.0), Inches(1.8), [
    ("Result: ", "+2.08 pp mean PG across 3 seeds"),
    ("Variance: ", "σ 2.57% → 0.41%  (6× reduction)"),
    ("Together: ", "v5 → v10  from 18.68% to 20.76% mean"),
], size=14, dot=ACCENT2)

add_text(sl, Inches(0.8), Inches(6.75), Inches(11.7), Inches(0.45),
         "Ablations varying loss weight, LoRA targets, and aggregation all returned to the same mean once both fixes were in place.",
         size=13, color=MED_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8  Variance collapse table
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Text-Path Fix: Variance Collapse", accent_w=Inches(5.2))

add_table(sl, Inches(0.8), Inches(1.35), Inches(11.7), Inches(1.9),
    ["Recipe", "Seed 0", "Seed 1", "Seed 2", "Mean ± Std", "Δ vs B0"],
    [
        ["v5  (mean tokens)",   "19.26%", "15.87%", "20.90%",
         "18.68 ± 2.57%",  ("+3.94 pp", AMBER)],
        [("v10  (EOS fix)", ACCENT), ("20.57%", GREEN), ("20.48%", GREEN), ("21.23%", GREEN),
         ("20.76 ± 0.41%", GREEN), ("+6.02 pp", GREEN)],
    ],
)

add_text(sl, Inches(0.8), Inches(3.55), Inches(11.7), Inches(0.55),
         "σ: 2.57%  ──────────────────────→  0.41%    (6× reduction)",
         size=26, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

add_bullets(sl, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.8), [
    ("Every v10 seed lands in [20.48, 21.23] ",
     "— a 0.75 pp spread vs. v5's 5.03 pp spread"),
    ("Best-PG snapshot: ",
     "peak PG = final PG on every seed — late-training drift eliminated"),
    ("Mechanistic reading: ",
     "EOS fix raised the peak by aligning training with evaluation. Snapshot locked it in."),
    ("R@1 halfmax: ",
     "7.30 ± 0.00% across all seeds — the model learns a stable broad signal; variance is in which patch wins argmax"),
], size=15)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9  Ablations — what didn't work
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "What Didn't Work: Ablations", accent_w=Inches(4.8))

cards = [
    ("↑ Region Loss Weight", "λ: 0.5 → 1.0", [
        "PG dropped −3.72 pp",
        "R@1 improved — broader, not sharper response",
        "Global loss provides structural regularisation",
        "At λ=1.0: collapse to PG < B0 at step 400",
    ], RED),
    ("↑ LoRA Capacity", "q+v → q+k+v  (+50% params)", [
        "Highest peak ever: 22.69% at step 400",
        "Drifted to 17.80% by end of epoch",
        "More capacity = more freedom to walk away",
        "Peak improved; final result did not",
    ], RED),
    ("Top-K Mean Loss", "max-pool → mean of top-3 patches", [
        "19.35% vs 19.96% — within noise",
        "R@1 halfmax: 7.30% identical on every variant",
        "Aggregation is not the lever",
        "Representation alignment is",
    ], RED),
]
for i, (title, sub, buls, col) in enumerate(cards):
    x = Inches(0.8 + i * 4.15)
    add_rect(sl, x, Inches(1.5), Inches(3.85), Inches(4.7),
             BG_CARD, border_color=col, corner_radius=0.03)
    add_text(sl, x + Inches(0.2), Inches(1.68), Inches(3.45), Inches(0.38),
             title, size=17, color=col, bold=True)
    add_text(sl, x + Inches(0.2), Inches(2.1), Inches(3.45), Inches(0.32),
             sub, size=13, color=LIGHT_GRAY)
    add_bullets(sl, x + Inches(0.2), Inches(2.5), Inches(3.45), Inches(3.3),
                buls, size=13, dot=col)

add_rect(sl, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.75),
         BG_CARD, border_color=ACCENT, corner_radius=0.03)
add_text(sl, Inches(1.1), Inches(6.5), Inches(11.1), Inches(0.55),
         "Takeaway: The bottleneck is representation consistency — not loss formulation or model capacity.",
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10  Qualitative hero — pg_viz_04  (1416×1479, ratio 0.957)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
badge(sl, "QUALITATIVE RESULTS", color=GREEN, w=Inches(3.3))

add_text(sl, Inches(0.8), Inches(1.05), Inches(12.0), Inches(0.55),
         "Most Dramatic Case — Image 1507563902  (n=3 phrases, each row = one phrase)",
         size=22, color=WHITE, bold=True)

# image: height=5.65" → width = 5.65 * 0.957 ≈ 5.41"
sl.shapes.add_picture(IMG_VIZ_04, Inches(0.45), Inches(1.62),
                      height=Inches(5.65))   # width auto ≈ 5.41"

# stats card  (starts at ~6.1" from left)
SX = Inches(6.3)
add_rect(sl, SX, Inches(1.62), Inches(6.65), Inches(5.65),
         BG_CARD, border_color=GREEN, corner_radius=0.03)

add_text(sl, SX + Inches(0.3), Inches(1.85), Inches(6.0), Inches(0.42),
         "Per-Model Pointing Game  (3 phrases)", size=17, color=ACCENT, bold=True)
add_line(sl, SX + Inches(0.3), Inches(2.33), Inches(5.9), color=BORDER)

scores = [("Frozen B0:", "33%   ✓✗✗", RED),
          ("M_human:",   "100%  ✓✓✓", GREEN),
          ("M_auto:",    "67%   ✓✓✗", AMBER)]
for j, (lbl, sc, col) in enumerate(scores):
    y = Inches(2.55 + j * 0.95)
    add_text(sl, SX + Inches(0.3), y, Inches(2.3), Inches(0.65),
             lbl, size=18, color=LIGHT_GRAY)
    add_text(sl, SX + Inches(2.65), y, Inches(3.8), Inches(0.65),
             sc, size=22, color=col, bold=True)

add_line(sl, SX + Inches(0.3), Inches(5.55), Inches(5.9), color=BORDER)
add_text(sl, SX + Inches(0.3), Inches(5.7), Inches(6.0), Inches(0.45),
         "B0 misses 2 of 3 phrases.  M_human hits all 3.", size=15, color=WHITE, bold=True)
add_text(sl, SX + Inches(0.3), Inches(6.25), Inches(6.0), Inches(0.75),
         "Plasma heatmap = cosine similarity of each patch to the query phrase\n"
         "★ = argmax patch  (green=HIT / red=MISS)   ■ = ground-truth bbox",
         size=13, color=MED_GRAY)

add_text(sl, Inches(0.45), Inches(7.3), Inches(5.5), Inches(0.25),
         "Columns: Frozen B0  |  M_human  |  M_auto", size=12, color=MED_GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11  Qualitative — pg_viz_01  (1345×1516, ratio 0.887)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
badge(sl, "QUALITATIVE RESULTS", color=GREEN, w=Inches(3.3))

add_text(sl, Inches(0.8), Inches(1.05), Inches(12.0), Inches(0.55),
         "Largest Combined Gain — Image 1836335410  (n=3 phrases)",
         size=22, color=WHITE, bold=True)

# image: height=5.65" → width = 5.65 * 0.887 ≈ 5.01"
sl.shapes.add_picture(IMG_VIZ_01, Inches(0.45), Inches(1.62),
                      height=Inches(5.65))

SX2 = Inches(6.1)
add_rect(sl, SX2, Inches(1.62), Inches(6.85), Inches(5.65),
         BG_CARD, border_color=GREEN, corner_radius=0.03)

add_text(sl, SX2 + Inches(0.3), Inches(1.85), Inches(6.2), Inches(0.42),
         "Per-Model Pointing Game  (3 phrases)", size=17, color=ACCENT, bold=True)
add_line(sl, SX2 + Inches(0.3), Inches(2.33), Inches(6.1), color=BORDER)

scores2 = [("Frozen B0:", "0%    ✗✗✗", RED),
           ("M_human:",   "100%  ✓✓✓", GREEN),
           ("M_auto:",    "33%   ✓✗✗", AMBER)]
for j, (lbl, sc, col) in enumerate(scores2):
    y = Inches(2.55 + j * 0.95)
    add_text(sl, SX2 + Inches(0.3), y, Inches(2.3), Inches(0.65),
             lbl, size=18, color=LIGHT_GRAY)
    add_text(sl, SX2 + Inches(2.65), y, Inches(3.8), Inches(0.65),
             sc, size=22, color=col, bold=True)

add_line(sl, SX2 + Inches(0.3), Inches(5.55), Inches(6.1), color=BORDER)
add_text(sl, SX2 + Inches(0.3), Inches(5.7), Inches(6.2), Inches(0.45),
         "Combined gain over B0: +133 pp across 3 phrases", size=15, color=WHITE, bold=True)
add_text(sl, SX2 + Inches(0.3), Inches(6.2), Inches(6.2), Inches(0.75),
         "Frozen model misses every phrase.\n"
         "Fine-tuned model (M_human) hits every phrase.", size=14, color=LIGHT_GRAY)

add_text(sl, Inches(0.45), Inches(7.3), Inches(5.5), Inches(0.25),
         "Columns: Frozen B0  |  M_human  |  M_auto", size=12, color=MED_GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12  Staircase — 5 epoch training
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Longer Training: Staircase Improvement", accent_w=Inches(6))

add_table(sl, Inches(0.8), Inches(1.35), Inches(6.5), Inches(3.7),
    ["Step", "Ep", "Val PG", "Test PG", "Behaviour"],
    [
        ["800",  "1", "19.68%", "19.92%", "Matches 1-ep result"],
        ["1600", "2", "19.40%", "20.74%", "Flat, oscillating"],
        ["2400", "3", "19.21%", "21.41%", "Still flat"],
        [("3000", AMBER), ("4", AMBER), ("23.59%", AMBER), ("24.59%", AMBER), ("STEP CHANGE +3 pp", AMBER)],
        [("3400 ★", GREEN), ("5", GREEN), ("24.44%", GREEN), ("25.26%", GREEN), ("Best-PG snapshot", GREEN)],
        ["4000",  "5", "24.11%", "26.03%", "Late epoch 5"],
    ],
    col_widths=[Inches(0.9), Inches(0.5), Inches(1.1), Inches(1.1), Inches(2.8)]
)

add_rect(sl, Inches(7.6), Inches(1.35), Inches(5.1), Inches(4.1),
         BG_CARD, border_color=GREEN, corner_radius=0.03)
add_text(sl, Inches(7.9), Inches(1.55), Inches(4.5), Inches(0.42),
         "Best Result (seed 1)", size=20, color=GREEN, bold=True)
add_text(sl, Inches(7.9), Inches(2.1), Inches(4.5), Inches(0.48),
         "Val PG = 24.44%", size=26, color=WHITE, bold=True)
add_text(sl, Inches(7.9), Inches(2.62), Inches(4.5), Inches(0.38),
         "+9.70 pp over B0", size=17, color=GREEN)
add_text(sl, Inches(7.9), Inches(3.1), Inches(4.5), Inches(0.48),
         "Test PG = 26.16%", size=26, color=WHITE, bold=True)
add_text(sl, Inches(7.9), Inches(3.62), Inches(4.5), Inches(0.38),
         "+11.42 pp over B0  (~48σ signal)", size=17, color=GREEN)
add_text(sl, Inches(7.9), Inches(4.12), Inches(4.5), Inches(0.38),
         "R@1 top-25 = 8.29%", size=17, color=AMBER)
add_text(sl, Inches(7.9), Inches(4.55), Inches(4.5), Inches(0.38),
         "First model to beat B0's 7.91%", size=14, color=LIGHT_GRAY)

add_rect(sl, Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.3),
         BG_CARD, border_color=AMBER, corner_radius=0.03)
add_text(sl, Inches(1.1), Inches(5.45), Inches(11.2), Inches(1.1),
         "Per-epoch cosine restarts kick the optimizer out of the epoch-1 basin. "
         "Epochs 1–3 reproduce the 1-epoch result; a step change occurs in epoch 4 (+3–5 pp). "
         "When ablations all fail, more training may be more productive than further hyperparameter search.",
         size=14, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13  Auto-annotations — vocabulary controls everything
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Auto-Annotations: Vocabulary Controls Everything", size=28, accent_w=Inches(7))

# Florence-2 viz strip  (3698×735, ratio 5.03)  width=12" → height≈2.39"
sl.shapes.add_picture(IMG_FLORENCE, Inches(0.67), Inches(1.18), width=Inches(12.0))
add_text(sl, Inches(0.67), Inches(3.63), Inches(12.0), Inches(0.3),
         "Florence-2 CAPTION_TO_PHRASE_GROUNDING — short noun-phrase → bbox annotations on Flickr30k images",
         size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(0.8), Inches(4.05), Inches(5.7), Inches(2.7),
         BG_CARD, border_color=RED, corner_radius=0.03)
add_text(sl, Inches(1.1), Inches(4.22), Inches(5.2), Inches(0.38),
         "M_auto: Dense Region Captions  ✗", size=17, color=RED, bold=True)
add_bullets(sl, Inches(1.1), Inches(4.65), Inches(5.2), Inches(2.0), [
    ("Labels: ", '"a woman in a blue top standing near a fence"'),
    ("PG ≈ 15.1% ", "— barely above B0 (14.74%)"),
    ("Root cause: ", "training vocabulary ≠ evaluation phrases"),
], size=14, dot=RED)

add_rect(sl, Inches(7.0), Inches(4.05), Inches(5.7), Inches(2.7),
         BG_CARD, border_color=GREEN, corner_radius=0.03)
add_text(sl, Inches(7.3), Inches(4.22), Inches(5.2), Inches(0.38),
         "M_cpg: Caption-to-Phrase Grounding  ✓", size=17, color=GREEN, bold=True)
add_bullets(sl, Inches(7.3), Inches(4.65), Inches(5.2), Inches(2.0), [
    ("Labels: ", '"woman",  "blue top",  "fence"'),
    ("PG = 23.35% ", "— +8.25 pp over M_auto"),
    ("Nearly matches M_human: ", "23.35% vs 24.44%"),
], size=14, dot=GREEN)

add_text(sl, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.42),
         "Same Florence-2 model  ·  Same images  ·  Same recipe  ·  Only the labels differ",
         size=16, color=AMBER, bold=True, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14  Main results table
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Main Results", accent_w=Inches(3.0))

add_table(sl, Inches(0.5), Inches(1.35), Inches(12.3), Inches(4.2),
    ["Model", "Config", "Val PG", "Test PG", "R@1 top-25", "R@1 halfmax"],
    [
        ["B0 (frozen)",         "—",                  "14.74%",              "—",                   "7.91%",              "7.58%"],
        ["M_human v5",          "FILIP, 1 ep",        "18.68 ± 2.57%",       "—",                   "7.12 ± 0.12%",       ("7.30 ± 0.00%", AMBER)],
        [("M_human v10", ACCENT),"EOS+snap, 1 ep",   ("20.76 ± 0.41%", ACCENT),("20.07 ± 1.67%", ACCENT),"—",            "7.28%"],
        [("M_human v11", GREEN), "EOS+snap, 5 ep",   ("24.44%", GREEN),     ("26.16%", GREEN),     ("8.29%", GREEN),     "7.30%"],
        [("M_auto", RED),        "Dense cap., 1 ep", ("~15.1%", RED),        "—",                   "—",                  "—"],
        [("M_cpg", GREEN),       "Phrase gnd., 5 ep",("23.35%", GREEN),     ("23.15%", GREEN),      "—",                  "—"],
    ],
    col_widths=[Inches(1.9), Inches(2.0), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1)]
)

add_text(sl, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.45),
         "Best single-seed: M_human v11 — 24.44% val / 26.16% test  (+9.70 / +11.42 pp over B0)",
         size=16, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.38),
         "M_cpg nearly matches M_human (23.35% vs 24.44%) using only Florence-2 phrase labels — no human bboxes needed",
         size=14, color=AMBER, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.8), Inches(6.75), Inches(11.7), Inches(0.32),
         "† Single seed; 3-seed mean ± std pending for v11 and M_cpg",
         size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15  Key takeaways
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, BG_DARK)
slide_title(sl, "Key Takeaways", size=36, accent_w=Inches(3.0))

takeaways = [
    ("Fix mismatches first. ",
     "Train/eval alignment matters more than loss design. Two one-line fixes account for the bulk of improvement."),
    ("Measure variance before claiming improvement. ",
     "The first single-seed result (22.08%) was +1.3σ noise. 3-seed verification is mandatory."),
    ("Longer training with restarts finds better basins. ",
     "Flat for 3 epochs, then a step change at epoch 4. When ablations fail — train longer."),
    ("For auto-labels, vocabulary is everything. ",
     "Accurate bboxes + wrong phrase distribution = no improvement. M_cpg nearly matches M_human."),
    ("295K params, 40 min on an L4. ",
     "Enough to teach SigLIP where phrases are — and composable with any existing SigLIP pipeline."),
]
for i, (bold_text, normal_text) in enumerate(takeaways):
    y = Inches(1.45 + i * 1.1)
    add_rect(sl, Inches(0.8), y, Inches(0.48), Inches(0.48), ACCENT, corner_radius=0.5)
    add_text(sl, Inches(0.8), y + Pt(2), Inches(0.48), Inches(0.48),
             str(i + 1), size=18, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    bold_normal_row(sl, Inches(1.5), y, Inches(11.0), Inches(0.85),
                    bold_text, normal_text, size=17)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16  Thank you
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
set_slide_bg(sl, RGBColor(0x0A, 0x10, 0x20))
add_text(sl, Inches(1), Inches(1.8), Inches(11), Inches(1),
         "Thank You!", size=54, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_line(sl, Inches(5), Inches(3.2), Inches(3.3))
add_text(sl, Inches(1), Inches(3.6), Inches(11), Inches(0.5),
         "Questions?", size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
         "github.com/Ani0202/region-grounded",
         size=18, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(5.35), Inches(11), Inches(0.5),
         "Siddharth Raj  ·  Aniket Agrawal",
         size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(5.9), Inches(11), Inches(0.5),
         "University of Chicago", size=16, color=MED_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
output_path = os.path.join(SCRIPT_DIR, "graft_presentation.pptx")
prs.save(output_path)
print(f"✓ Saved: {output_path}")
print(f"  Slides: {len(prs.slides)}")
